import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

tables_data = [
    ("Categories", [
        ("Field", "Type", "Key", "Description"),
        ("cid", "INT", "PK", "Unique category ID"),
        ("cname", "VARCHAR", "", "Category name"),
        ("description", "VARCHAR", "", "Category description")
    ]),
    ("Complaints", [
        ("Field", "Type", "Key", "Description"),
        ("complaint_id", "INT", "PK", "Unique complaint ID"),
        ("uid", "INT", "FK", "User reference"),
        ("subject", "VARCHAR", "", "Complaint subject"),
        ("description", "TEXT", "", "Complaint details"),
        ("status", "ENUM", "", "Complaint status"),
        ("created_at", "TIMESTAMP", "", "Creation time"),
        ("resolved_at", "TIMESTAMP", "", "Resolution time")
    ]),
    ("Image Table", [
        ("Field", "Type", "Key", "Description"),
        ("image_id", "INT", "PK", "Unique image ID"),
        ("pid", "INT", "FK", "Product reference"),
        ("image_url", "VARCHAR", "", "Image URL"),
        ("is_primary", "TINYINT", "", "Primary image flag"),
        ("uploaded_at", "TIMESTAMP", "", "Upload time")
    ]),
    ("Payments", [
        ("Field", "Type", "Key", "Description"),
        ("payment_id", "INT", "PK", "Unique payment ID"),
        ("purchase_id", "INT", "FK", "Purchase reference"),
        ("rental_id", "INT", "FK", "Rental reference"),
        ("payer_id", "INT", "FK", "User (payer) reference"),
        ("amount", "DECIMAL", "", "Payment amount"),
        ("payment_method", "ENUM", "", "Method of payment"),
        ("payment_status", "ENUM", "", "Status of payment"),
        ("transaction_ref", "VARCHAR", "", "Transaction reference"),
        ("payment_date", "TIMESTAMP", "", "Date of payment")
    ]),
    ("Product Rent", [
        ("Field", "Type", "Key", "Description"),
        ("prid", "INT", "PK", "Unique product rent ID"),
        ("pid", "INT", "FK", "Product reference"),
        ("no_of_days", "INT", "", "Number of days"),
        ("charge_per_day", "DECIMAL", "", "Rent charge per day"),
        ("security_deposit", "DECIMAL", "", "Security deposit")
    ]),
    ("Products", [
        ("Field", "Type", "Key", "Description"),
        ("pid", "INT", "PK", "Unique product ID"),
        ("uid", "INT", "FK", "User (owner) reference"),
        ("cid", "INT", "FK", "Category reference"),
        ("pname", "VARCHAR", "", "Product name"),
        ("pdesc", "TEXT", "", "Product description"),
        ("price", "DECIMAL", "", "Product price"),
        ("status", "ENUM", "", "Product status"),
        ("type", "ENUM", "", "Type (SELL, RENT, BOTH)"),
        ("created_at", "TIMESTAMP", "", "Creation time")
    ]),
    ("Purchases", [
        ("Field", "Type", "Key", "Description"),
        ("purchase_id", "INT", "PK", "Unique purchase ID"),
        ("pid", "INT", "FK", "Product reference"),
        ("buyer_id", "INT", "FK", "User (buyer) reference"),
        ("seller_id", "INT", "FK", "User (seller) reference"),
        ("purchase_date", "DATETIME", "", "Date of purchase"),
        ("amount", "DECIMAL", "", "Purchase amount"),
        ("status", "ENUM", "", "Purchase status")
    ]),
    ("Rental Transactions", [
        ("Field", "Type", "Key", "Description"),
        ("rental_id", "INT", "PK", "Unique rental ID"),
        ("pid", "INT", "FK", "Product reference"),
        ("owner_id", "INT", "FK", "User (owner) reference"),
        ("renter_id", "INT", "FK", "User (renter) reference"),
        ("start_date", "DATE", "", "Start date"),
        ("end_date", "DATE", "", "End date"),
        ("total_amount", "DECIMAL", "", "Total amount"),
        ("status", "ENUM", "", "Rental status")
    ]),
    ("Reports", [
        ("Field", "Type", "Key", "Description"),
        ("report_id", "INT", "PK", "Unique report ID"),
        ("pid", "INT", "FK", "Product reference"),
        ("reporter_id", "INT", "FK", "User (reporter) reference"),
        ("reason", "VARCHAR", "", "Report reason"),
        ("description", "TEXT", "", "Report details"),
        ("status", "ENUM", "", "Report status"),
        ("report_date", "TIMESTAMP", "", "Date of report")
    ]),
    ("Reviews", [
        ("Field", "Type", "Key", "Description"),
        ("review_id", "INT", "PK", "Unique review ID"),
        ("pid", "INT", "FK", "Product reference"),
        ("uid", "INT", "FK", "User reference"),
        ("rating", "INT", "", "Rating (1-5)"),
        ("review_text", "TEXT", "", "Review details"),
        ("review_date", "TIMESTAMP", "", "Date of review")
    ]),
    ("Roles", [
        ("Field", "Type", "Key", "Description"),
        ("rid", "INT", "PK", "Unique role ID"),
        ("rname", "VARCHAR", "", "Role name")
    ]),
    ("Users", [
        ("Field", "Type", "Key", "Description"),
        ("uid", "INT", "PK", "Unique user ID"),
        ("uname", "VARCHAR", "", "Username"),
        ("password", "VARCHAR", "", "User password"),
        ("email", "VARCHAR", "", "Email address"),
        ("contactnumber", "VARCHAR", "", "Contact number"),
        ("fname", "VARCHAR", "", "First name"),
        ("lname", "VARCHAR", "", "Last name"),
        ("rid", "INT", "FK", "Role reference"),
        ("created_at", "TIMESTAMP", "", "Account creation time"),
        ("status", "ENUM", "", "User status")
    ])
]

def set_cell_border(cell):
    from pptx.oxml.xmlchemy import OxmlElement
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = OxmlElement(border)
        ln.set('w', '12700')
        ln.set('cmpd', 'sng')
        solidFill = OxmlElement('a:solidFill')
        srgbClr = OxmlElement('a:srgbClr')
        srgbClr.set('val', '000000')
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)

def set_cell_text(cell, text, is_bold=False):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = 'Times New Roman'
            run.font.size = Pt(11)
            run.font.bold = is_bold
            run.font.color.rgb = RGBColor(0, 0, 0)
    
    set_cell_border(cell)

prs = Presentation()
blank_slide_layout = prs.slide_layouts[5] # 5 is title only, 6 is blank

slide = None
y_offset = Inches(0.5)

margin_left = Inches(1)
table_width = Inches(8)
col_widths = [Inches(2.0), Inches(1.5), Inches(1.0), Inches(3.5)]

for i, (table_name, rows) in enumerate(tables_data):
    if i % 3 == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        y_offset = Inches(0.5)
        
        if i == 0:
            # Title for first slide
            txBox = slide.shapes.add_textbox(margin_left, y_offset, Inches(8), Inches(0.5))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "7. Database Design:"
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = 'Times New Roman'
            y_offset += Inches(0.8)

    # Add table title
    txBox = slide.shapes.add_textbox(margin_left, y_offset, Inches(8), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = f"{i+1}.  {table_name}"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.name = 'Times New Roman'
    
    y_offset += Inches(0.4)
    
    # Add table
    rows_count = len(rows)
    cols_count = len(rows[0])
    
    table_shape = slide.shapes.add_table(rows_count, cols_count, margin_left, y_offset, table_width, Inches(0.3 * rows_count))
    table = table_shape.table
    
    for c in range(cols_count):
        table.columns[c].width = col_widths[c]
    
    for r_idx, row in enumerate(rows):
        for c_idx, cell_value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            set_cell_text(cell, cell_value, is_bold=(r_idx==0))
            # clear cell background fill
            cell.fill.background()
            
    y_offset += Inches(0.3 * rows_count) + Inches(0.5)

prs.save('Database_Design.pptx')
print("PPTX generated successfully!")
